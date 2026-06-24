import os
import shutil
import subprocess
import tempfile
import mailbox
from email import policy
from email.parser import BytesParser

from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
import fitz
import pandas as pd
import docx2txt
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

from file_utilities import FileFormats, get_nonexistant_path


def return_libreoffice_path():
    return shutil.which("soffice") or shutil.which("libreoffice")


def external_dependencies_check():
    print("Checking Linux dependencies...")
    print(f"LibreOffice: {return_libreoffice_path() or 'NOT FOUND'}")
    print(f"Tesseract: {shutil.which('tesseract') or 'NOT FOUND'}")
    print(f"FFmpeg: {shutil.which('ffmpeg') or 'NOT FOUND'}")
    print("PST/OST: disabled in Linux port unless a native extractor is added")
    print("Dependency check complete.")


def _write(text, src_path, output_folder):
    text = text or ""
    if output_folder and os.path.isdir(output_folder):
        out = get_nonexistant_path(os.path.join(output_folder, os.path.basename(src_path) + ".extracted.txt"))
        with open(out, "w", encoding="utf-8", errors="replace") as f:
            f.write(text)
        print(f"Wrote extracted text to {out}")
    return [text]


def process_html(html):
    return BeautifulSoup(html or "", "html.parser").get_text("\n")


def extract_txt(path, output_folder=None, ocr=False):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return _write(f.read(), path, output_folder)


def extract_rtf(path, output_folder=None, ocr=False):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return _write(rtf_to_text(f.read()), path, output_folder)


def extract_html(path, output_folder=None, ocr=False):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return _write(process_html(f.read()), path, output_folder)


def extract_docx(path, output_folder=None, ocr=False):
    return _write(docx2txt.process(path) or "", path, output_folder)


def extract_pptx(path, output_folder=None, ocr=False):
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        chunks.append(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
        if getattr(slide, "has_notes_slide", False):
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                chunks.append("\n[Notes]\n" + notes)
    return _write("\n".join(chunks), path, output_folder)


def extract_excel(path, output_folder=None, ocr=False):
    chunks = []
    sheets = pd.read_excel(path, sheet_name=None, dtype=str)
    for sheet_name, df in sheets.items():
        chunks.append(f"\n--- Sheet: {sheet_name} ---\n")
        chunks.append(df.fillna("").to_csv(index=False))
    return _write("\n".join(chunks), path, output_folder)


def extract_pdf(path, output_folder=None, ocr=False):
    chunks = []
    doc = fitz.open(path)
    for page_num, page in enumerate(doc, 1):
        chunks.append(f"\n--- Page {page_num} ---\n")
        txt = page.get_text() or ""
        chunks.append(txt)

        if ocr and not txt.strip():
            pix = page.get_pixmap(dpi=200)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(pix.tobytes("png"))
                tmp_path = tmp.name
            try:
                chunks.append(pytesseract.image_to_string(Image.open(tmp_path)))
            finally:
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

    return _write("\n".join(chunks), path, output_folder)


def extract_image(path, output_folder=None, ocr=True):
    text = pytesseract.image_to_string(Image.open(path)) if ocr else ""
    return _write(text, path, output_folder)


def _convert_with_libreoffice(path, output_folder, target):
    lo = return_libreoffice_path()
    if not lo:
        return [f"LibreOffice not found; cannot convert {path}"]

    with tempfile.TemporaryDirectory() as td:
        result = subprocess.run(
            [lo, "--headless", "--convert-to", target, "--outdir", td, path],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            return [f"LibreOffice conversion failed:\n{result.stderr}"]

        converted = None
        for name in os.listdir(td):
            if name.lower().endswith("." + target):
                converted = os.path.join(td, name)
                break

        if not converted:
            return [f"LibreOffice did not produce .{target}"]

        if target == "docx":
            return extract_docx(converted, output_folder)
        if target == "pptx":
            return extract_pptx(converted, output_folder)

    return ["Unsupported conversion"]


def extract_doc(path, output_folder=None, ocr=False):
    return _convert_with_libreoffice(path, output_folder, "docx")


def extract_ppt(path, output_folder=None, ocr=False):
    return _convert_with_libreoffice(path, output_folder, "pptx")


def extract_audio(path, output_folder=None, from_lang="Auto"):
    try:
        from faster_whisper import WhisperModel
    except Exception as e:
        return [f"Audio/video extraction disabled: Python faster-whisper is not installed.\n{e}"]

    outdir = output_folder if output_folder and os.path.isdir(output_folder) else tempfile.mkdtemp()

    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    model_dir = os.path.join(
        repo_root,
        "whisper_files",
        "Faster-Whisper-XXL",
        "_models",
        "faster-whisper-large-v2",
    )

    if not os.path.isdir(model_dir):
        return [f"Model folder not found:\n{model_dir}"]

    try:
        print(f"Loading model from {model_dir}")
        model = WhisperModel(model_dir, device="cpu", compute_type="int8")

        print(f"Transcribing {path}")
        segments, info = model.transcribe(path, beam_size=5)

        lines = []
        lines.append(f"Detected language: {info.language} ({info.language_probability:.2f})")
        lines.append("")

        for seg in segments:
            lines.append(f"[{seg.start:.2f} --> {seg.end:.2f}] {seg.text}")

        text = "\n".join(lines)

        if output_folder and os.path.isdir(output_folder):
            out = get_nonexistant_path(
                os.path.join(output_folder, os.path.basename(path) + ".transcript.txt")
            )
            with open(out, "w", encoding="utf-8", errors="replace") as f:
                f.write(text)
            print(f"Wrote transcript to {out}")

        return [text]

    except Exception as e:
        return [f"faster-whisper transcription failed:\n{e}"]

def extract_eml(path, output_folder=None, extract_attachments=False, process_attachments=False):
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    subject = msg.get("subject", "")
    headers = f"FROM: {msg.get('from','')}\nTO: {msg.get('to','')}\nDATE: {msg.get('date','')}\nSUBJECT: {subject}\n"
    text_body = ""
    html_body = ""

    for part in msg.walk():
        ctype = part.get_content_type()
        disp = part.get_content_disposition()

        if disp == "attachment":
            continue
        if ctype == "text/plain":
            try:
                text_body += part.get_content() + "\n"
            except Exception:
                pass
        elif ctype == "text/html":
            try:
                html_body += process_html(part.get_content()) + "\n"
            except Exception:
                pass

    final = headers + "\nTEXT BODY:\n" + text_body + "\nHTML BODY:\n" + html_body
    _write(final, path, output_folder)
    return [final, subject, text_body, html_body, headers]


def extract_mbox(path, output_folder=None, extract_attachments=False, process_attachments=False):
    chunks = []
    mbox = mailbox.mbox(path)
    for idx, msg in enumerate(mbox, 1):
        chunks.append(f"\n--- MESSAGE {idx} ---\n")
        chunks.append(f"FROM: {msg.get('from','')}\nTO: {msg.get('to','')}\nDATE: {msg.get('date','')}\nSUBJECT: {msg.get('subject','')}\n")
        payload = msg.get_payload(decode=True)
        if payload:
            chunks.append(payload.decode(errors="replace"))
    return _write("\n".join(chunks), path, output_folder)


def extract_pst(path, output_folder=None, extract_attachments=False, process_attachments=False):
    return ["PST/OST extraction is disabled in this Linux port. Original WHISPR used a Windows XstExport.exe helper."]


def extract_file(path, output_folder=None, file_format=FileFormats.AUTO, extract_attachments=False, process_attachments=False, ocr=False, from_lang="Auto"):
    if file_format == FileFormats.AUTO:
        return extract_file_autodetect(path, output_folder, extract_attachments, process_attachments, ocr, from_lang)

    if file_format == FileFormats.PDF:
        return extract_pdf(path, output_folder, ocr)
    if file_format == FileFormats.XLS:
        return extract_excel(path, output_folder, ocr)
    if file_format == FileFormats.DOCX:
        return extract_docx(path, output_folder, ocr)
    if file_format == FileFormats.DOC:
        return extract_doc(path, output_folder, ocr)
    if file_format == FileFormats.PPTX:
        return extract_pptx(path, output_folder, ocr)
    if file_format == FileFormats.PPT:
        return extract_ppt(path, output_folder, ocr)
    if file_format == FileFormats.RTF:
        return extract_rtf(path, output_folder, ocr)
    if file_format == FileFormats.TXT:
        return extract_txt(path, output_folder, ocr)
    if file_format == FileFormats.HTML:
        return extract_html(path, output_folder, ocr)
    if file_format == FileFormats.IMAGE:
        return extract_image(path, output_folder, ocr)
    if file_format == FileFormats.AV:
        return extract_audio(path, output_folder, from_lang)
    if file_format == FileFormats.EML:
        return extract_eml(path, output_folder, extract_attachments, process_attachments)
    if file_format == FileFormats.MBOX:
        return extract_mbox(path, output_folder, extract_attachments, process_attachments)
    if file_format == FileFormats.PST:
        return extract_pst(path, output_folder, extract_attachments, process_attachments)

    return [f"Unsupported file format: {file_format}"]


def extract_file_autodetect(path, output_folder=None, extract_attachments=False, process_attachments=False, ocr=False, from_lang="Auto"):
    ext = os.path.splitext(path)[1].lower().lstrip(".")

    if ext == "pdf":
        return extract_pdf(path, output_folder, ocr)
    if ext in ("xlsx", "xlsm", "xls"):
        return extract_excel(path, output_folder, ocr)
    if ext in ("docx", "docm"):
        return extract_docx(path, output_folder, ocr)
    if ext == "doc":
        return extract_doc(path, output_folder, ocr)
    if ext in ("pptx", "pptm"):
        return extract_pptx(path, output_folder, ocr)
    if ext == "ppt":
        return extract_ppt(path, output_folder, ocr)
    if ext == "rtf":
        return extract_rtf(path, output_folder, ocr)
    if ext in ("txt", "text", "csv", "log"):
        return extract_txt(path, output_folder, ocr)
    if ext in ("html", "htm"):
        return extract_html(path, output_folder, ocr)
    if ext in ("png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"):
        return extract_image(path, output_folder, ocr)
    if ext in ("mp3", "wav", "m4a", "mp4", "mkv", "avi", "mov", "flac", "ogg"):
        return extract_audio(path, output_folder, from_lang)
    if ext == "eml":
        return extract_eml(path, output_folder, extract_attachments, process_attachments)
    if ext == "mbox":
        return extract_mbox(path, output_folder, extract_attachments, process_attachments)
    if ext in ("pst", "ost"):
        return extract_pst(path, output_folder, extract_attachments, process_attachments)

    return [f"Unsupported or unknown file type: {path}"]
